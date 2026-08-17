"""
OpsTracker - SPF NS Reservist Attendance Management
Approval deck for commanding officer / planning committee.
No em dashes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ---- PALETTE ----
NAVY    = RGBColor(0x00, 0x20, 0x5B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x4F, 0x81, 0xBD)
LGRAY   = RGBColor(0xD4, 0xE1, 0xF2)
GREEN   = RGBColor(0x9B, 0xBB, 0x59)
RED_T   = RGBColor(0xFC, 0xE4, 0xD6)
GRN_T   = RGBColor(0xE2, 0xEF, 0xDA)
DK_RED  = RGBColor(0x7F, 0x1F, 0x1F)
DK_GRN  = RGBColor(0x1F, 0x4E, 0x1F)
GRN_PIL = RGBColor(0x26, 0x6B, 0x4A)
AMBER   = RGBColor(0xE0, 0x8A, 0x00)

# Content area from Layout 3 placeholder
BX, BY, BW, BH = 1.78, 1.15, 10.87, 5.80
# Precomputed column widths
CW2 = (BW - 0.20) / 2      # two-column: 5.335" each
CW3 = (BW - 0.40) / 3      # three-column: 3.49" each
# Row heights for 2x2 section layouts
# Slides with footer note (note at y=6.45): fill sections to y=6.35
RH_N, RG_N = 2.52, 0.15    # row height, gap
BOT_Y_N = BY + RH_N + RG_N  # top of bottom row (=3.82)
# Slides without footer note: fill sections to y=6.85
RH, RG = 2.77, 0.15
BOT_Y = BY + RH + RG        # top of bottom row (=4.07)

# ---- HELPERS ----

def remove_all_slides(prs):
    ids = prs.slides._sldIdLst
    for elem in list(ids):
        rId = elem.get(qn('r:id'))
        ids.remove(elem)
        prs.part.drop_rel(rId)

def _tf(slide, l, t, w, h):
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    sh.text_frame.word_wrap = True
    return sh.text_frame

def _p(tf, text, sz=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
       sp=6, first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(sp)
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = 'Arial'
    return p

def _pill(slide, l, t, w, h, label, bg, fg=WHITE, sz=13):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = bg; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(sz); run.font.bold = True
    run.font.color.rgb = fg; run.font.name = 'Arial'

def _bullets(tf, items, sz=14, sp0=8, sp=11, color=NAVY):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(sp0 if i == 0 else sp)
        run = p.add_run(); run.text = '•  ' + item
        run.font.size = Pt(sz); run.font.color.rgb = color; run.font.name = 'Arial'

def _section_col(slide, l, t, w, h, label, bg, items, sz=14):
    """Card section: LGRAY background, colored pill header, bullet list."""
    ph = 0.40
    # Card background
    card = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = LGRAY; card.line.fill.background()
    # Colored pill header
    _pill(slide, l, t, w, ph, label, bg)
    # Bullet content (indented slightly inside card)
    tf = _tf(slide, l + 0.10, t + ph + 0.10, w - 0.20, h - ph - 0.15)
    _bullets(tf, items, sz=sz, sp0=6, sp=10, color=NAVY)

def _cell(cell, text, bg=None, fg=NAVY, sz=12, bold=False, align=PP_ALIGN.LEFT):
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell._tc.get_or_add_tcPr().set('anchor', 'ctr')
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for r in list(p.runs): p._p.remove(r._r)
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.color.rgb = fg; run.font.name = 'Arial'

def _note(slide, text, y=6.45, sz=10):
    tf = _tf(slide, BX, y, BW, 0.42)
    _p(tf, text, sz=sz, color=BLUE, italic=True, first=True)

def _pill_row(slide, items, t, h=0.44):
    n = len(items)
    pw = (BW - 0.15 * (n - 1)) / n
    for i, (lbl, bg) in enumerate(items):
        _pill(slide, BX + i * (pw + 0.15), t, pw, h, lbl, bg)

def _divider(slide, y, color=BLUE):
    d = slide.shapes.add_shape(1, Inches(BX), Inches(y), Inches(BW), Inches(0.025))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()

# ---- LOAD ----
prs = Presentation('OpsTracker.pptx')
L0 = prs.slide_layouts[0]
L3 = prs.slide_layouts[3]
remove_all_slides(prs)

# ============================================================
# 1.  TITLE
# ============================================================
s = prs.slides.add_slide(L0)
ph = s.placeholders[0]; ph.text = ''
tf0 = ph.text_frame
pm = tf0.paragraphs[0]; pm.alignment = PP_ALIGN.CENTER
rm = pm.add_run(); rm.text = 'OpsTracker'
rm.font.size = Pt(44); rm.font.bold = True
rm.font.color.rgb = WHITE; rm.font.name = 'Arial'
ps = tf0.add_paragraph(); ps.alignment = PP_ALIGN.CENTER; ps.space_before = Pt(8)
rs = ps.add_run(); rs.text = 'Digital Attendance Management for SPF NS Reservists'
rs.font.size = Pt(20); rs.font.color.rgb = WHITE; rs.font.name = 'Arial'

s.placeholders[11].text = '17 August 2026'
s.placeholders[12].text = 'For Approval'
tf13 = s.placeholders[13].text_frame
p1 = tf13.paragraphs[0]; r1 = p1.add_run()
r1.text = 'Ops Security, Ops Branch'; r1.font.bold = True; r1.font.name = 'Arial'
p2 = tf13.add_paragraph(); r2 = p2.add_run()
r2.text = 'NSPI Sun Jin'; r2.font.name = 'Arial'

# ============================================================
# 2.  AGENDA
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Agenda'

agenda = [
    ('01', 'The Problem'),
    ('02', 'Before vs After'),
    ('03', 'Introducing OpsTracker'),
    ('04', 'Features and Capabilities'),
    ('05', 'Data Security and Privacy'),
    ('06', 'PDPA Compliance'),
    ('07', 'App Screenshots'),
    ('08', 'Expected Benefits and Impact'),
    ('09', 'Cost and Value'),
    ('10', 'Challenges'),
    ('11', 'Deployment Timeline'),
    ('12', 'Handover Plan'),
    ('13', 'Future Plans and Recommendation'),
]

def _agenda_col(slide, l, rows):
    tf = _tf(slide, l, BY, 4.5, 5.70)
    for i, (num, label) in enumerate(rows):
        _p(tf, num,   sz=22, bold=True,  color=BLUE, sp=14 if i > 0 else 4, first=(i == 0))
        _p(tf, label, sz=14, bold=False, color=NAVY, sp=1)

_agenda_col(s, BX,        agenda[:7])
_agenda_col(s, BX + 5.47, agenda[7:])

# ============================================================
# 3.  THE PROBLEM
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'The Problem'
tf = s.placeholders[12].text_frame; tf.word_wrap = True

_bullets(tf, [
    'Attendance based on trust alone - no system verification',
    'No way to confirm physical presence at reporting location',
    'Commanders have no real-time view of who has reported',
    'Leave and MC handled informally via chat with no record',
    'No consolidated attendance history exists for the unit',
    'Chronic absences go undetected until too late to act',
], sz=17, sp0=16, sp=18)

# ============================================================
# 4.  BEFORE vs AFTER
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Before vs After'

bva = [
    ('Attendance Check', 'Verbal, trust-based; no verification',
     'GPS-verified check-in via mobile browser'),
    ('Supervisor View',  'Must be on site to know who showed up',
     'Live dashboard updated in real time'),
    ('Leave and MC',     'WhatsApp or verbal; no formal record',
     'Digital submission with full approval history'),
    ('Record Keeping',   'Relies on memory and scattered chat logs',
     'Permanent digital log with one-click export'),
    ('Accountability',   'No mechanism to flag repeated no-shows',
     'Auto-absent marking and trend tracking'),
]

tbl = s.shapes.add_table(6, 3, Inches(BX - 0.18), Inches(BY), Inches(BW), Inches(BH)).table
for idx, w in enumerate([2.2, 4.15, 4.52]):
    tbl.columns[idx].width = Inches(w)

_cell(tbl.cell(0, 0), 'Area',       bg=NAVY, fg=WHITE, sz=13, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 1), 'Before',     bg=NAVY, fg=WHITE, sz=13, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 2), 'OpsTracker', bg=BLUE, fg=WHITE, sz=13, bold=True, align=PP_ALIGN.CENTER)

for ri, (area, before, after) in enumerate(bva):
    alt = ri % 2 == 0
    _cell(tbl.cell(ri+1, 0), area,   bg=LGRAY if alt else WHITE, fg=NAVY, sz=13, bold=True)
    _cell(tbl.cell(ri+1, 1), before, bg=RED_T, fg=DK_RED, sz=13)
    _cell(tbl.cell(ri+1, 2), after,  bg=GRN_T, fg=DK_GRN, sz=13)

# ============================================================
# 5.  INTRODUCING OPSTRACKER
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Introducing OpsTracker'

tf_i = _tf(s, BX, BY, BW, 0.42)
_p(tf_i,
   'Mobile-first Progressive Web App replacing verbal check-ins with GPS-verified digital attendance records.',
   sz=15, color=NAVY, first=True)

_pill_row(s, [
    ('Progressive Web App',    NAVY),
    ('Supabase + Realtime',    BLUE),
    ('GPS Shift Verification', GRN_PIL),
], t=1.68)

tf_p = _tf(s, BX, 2.24, BW, 4.66)
features = [
    ('Mobile-First  ',        'Works on iOS, Android, or desktop. No app store download required.'),
    ('GPS Check-In  ',        'Confirms device location before recording. Supports AM, PM, and Full Day shifts.'),
    ('Real-Time Dashboard  ', 'Live view of present, absent, and MC status as check-ins happen.'),
    ('Better Tracking  ',     'Full attendance history and trends for any reservist or cycle.'),
    ('Automated Reporting  ', 'One-click export to Excel or a print-ready cycle report.'),
]
for i, (lbl, det) in enumerate(features):
    p = tf_p.paragraphs[0] if i == 0 else tf_p.add_paragraph()
    p.space_before = Pt(9 if i > 0 else 2)
    r1 = p.add_run(); r1.text = lbl
    r1.font.size = Pt(15); r1.font.bold = True
    r1.font.color.rgb = BLUE; r1.font.name = 'Arial'
    r2 = p.add_run(); r2.text = det
    r2.font.size = Pt(15); r2.font.color.rgb = NAVY; r2.font.name = 'Arial'

# ============================================================
# 6.  RESERVIST EXPERIENCE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Reservist Experience'
tf = s.placeholders[12].text_frame; tf.word_wrap = True

steps = [
    ('Step 1  Register  ',    'Submit online form; admin activates your account.'),
    ('Step 2  Check In  ',    'Open app in reporting window; GPS confirms; one tap records attendance.'),
    ('Step 3  Leave / MC  ',  'Submit request digitally; receive approval or rejection in-app.'),
    ('Step 4  View History  ','Access personal attendance records and leave outcomes anytime.'),
]
for i, (lbl, det) in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(8 if i == 0 else 24)
    r1 = p.add_run(); r1.text = lbl
    r1.font.size = Pt(17); r1.font.bold = True
    r1.font.color.rgb = BLUE; r1.font.name = 'Arial'
    r2 = p.add_run(); r2.text = det
    r2.font.size = Pt(17); r2.font.color.rgb = NAVY; r2.font.name = 'Arial'

_note(s, 'Reservists can install OpsTracker to their home screen as a PWA. No app store required.')

# ============================================================
# 7.  ADMIN AND SUPERVISOR CAPABILITIES
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Admin and Supervisor Capabilities'

# 4 sections filling to y=6.35 (note slides), each row h=RH_N=2.52
_section_col(s, BX,          BY,       CW2, RH_N, 'Attendance', NAVY, [
    'Live dashboard: present, absent, and MC',
    'Manual attendance override with audit log',
    'Mark all present or set no-report days',
    'Auto-absent marking after check-in window closes',
], sz=14)
_section_col(s, BX,          BOT_Y_N,  CW2, RH_N, 'Leave and MC', BLUE, [
    'Approve or reject requests individually or in bulk',
    'Cancel approved leave if plans change',
    'Full leave history per reservist',
    'MC upload tracked separately from standard leave',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,       CW2, RH_N, 'Roster', NAVY, [
    'Add, deactivate, or reassign reservists',
    'CSV bulk upload and password resets',
    'Welfare notes and member activity log',
    'Assign reservists to batches and companies',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y_N,  CW2, RH_N, 'Reporting', BLUE, [
    'Export cycle report to Excel or print',
    'WhatsApp broadcast to the full unit',
    'Cycle picker and multi-day date range view',
    'Attendance trend snapshot per exercise cycle',
], sz=14)

_note(s, 'Super Admins can also manage admin accounts, create cycles, and define GPS reporting boundaries.')

# ============================================================
# 8.  DATA SECURITY AND PRIVACY
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Data Security and Privacy'

# 4 sections filling to y=6.85 (no note), each row h=RH=2.77
_section_col(s, BX,          BY,      CW2, RH, 'Access Control', NAVY, [
    'Three-tier roles: Reservist, Admin, Super Admin',
    'Each role restricted to authorised data only',
    'Super Admin approval required for account changes',
    'All actions tied to authenticated identity',
], sz=14)
_section_col(s, BX,          BOT_Y,   CW2, RH, 'Session Security', BLUE, [
    'Automatic logout after a period of inactivity',
    'Secure auth tokens managed by Supabase Auth',
    'No personal data stored on the device',
    'Login required on every new session',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,      CW2, RH, 'Data Protection', NAVY, [
    'PostgreSQL: encrypted at rest and in transit (TLS)',
    'Row-level security enforced at database layer',
    'No third-party data sharing',
    'ISO 27001 certified; SOC 2 Type II compliant',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y,   CW2, RH, 'Audit Trail', GRN_PIL, [
    'Every edit logged with timestamp and user',
    'Leave decisions and admin actions recorded',
    'Full history available for compliance review',
    'Immutable log - edits cannot be silently deleted',
], sz=14)

_note(s, 'Hosted externally on Supabase - comparable in approach to Google Forms, with stronger security controls. No SPF infrastructure required.')

# ============================================================
# 9.  PDPA COMPLIANCE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'PDPA Compliance'

tf_i = _tf(s, BX, BY, BW, 0.44)
_p(tf_i,
   'Designed to align with the Personal Data Protection Act 2012 (PDPA) and government data governance principles.',
   sz=14, color=NAVY, first=True)

# 4 sections from y=1.75 to y=6.85, each h=2.43
P_TOP  = BY + 0.60   # 1.75
P_RH   = 2.43
P_BOT  = P_TOP + P_RH + 0.17  # 4.35

_section_col(s, BX,          P_TOP,  CW2, P_RH, 'Data Collected', NAVY, [
    'Name and contact number (at registration)',
    'GPS location (at check-in only)',
    'Attendance and leave records',
    'No NRIC or financial data collected',
], sz=14)
_section_col(s, BX,          P_BOT,  CW2, P_RH, 'Individual Rights', BLUE, [
    'Reservists may request their own records',
    'Corrections made by authorised admin',
    'Data deleted after NS liability ends',
    'Requests handled within a reasonable time',
], sz=14)
_section_col(s, BX + CW2 + 0.20, P_TOP,  CW2, P_RH, 'Key Obligations Met', NAVY, [
    'Consent obtained at sign-up before data collected',
    'Purpose limited to attendance management only',
    'Only necessary data collected (data minimisation)',
    'Encryption, access control, row-level security',
], sz=14)
_section_col(s, BX + CW2 + 0.20, P_BOT,  CW2, P_RH, 'Recommended Actions', AMBER, [
    'Seek Data Protection Officer guidance before deployment',
    'Define data retention period per SPF policy',
    'Include a privacy notice at reservist sign-up',
    'Document the lawful basis for data collection',
], sz=14)

# ============================================================
# 10.  APP SCREENSHOTS
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'App Screenshots'

_note(s, '[ Replace the boxes below with actual screenshots before presenting ]', y=BY - 0.05)

SCR_W = (BW - 0.22) / 2   # 5.325
SCR_H = (BH - 0.60) / 2   # 2.60
SCR_T1 = BY + 0.52         # 1.67
SCR_T2 = SCR_T1 + SCR_H + 0.18  # 4.45
SCR_L2 = BX + SCR_W + 0.22      # 7.325

def _scrshot(slide, l, t, w, h, title, sub):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = LGRAY
    sh.line.color.rgb = BLUE; sh.line.width = Pt(1.5)
    tf = sh.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = NAVY; r1.font.name = 'Arial'
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(10); r2.font.italic = True; r2.font.color.rgb = BLUE; r2.font.name = 'Arial'

_scrshot(s, BX,    SCR_T1, SCR_W, SCR_H, '[Reservist: Check-In Screen]',  'GPS verification and shift phase')
_scrshot(s, SCR_L2, SCR_T1, SCR_W, SCR_H, '[Admin: Attendance Dashboard]', 'Live present, absent, and MC view')
_scrshot(s, BX,    SCR_T2, SCR_W, SCR_H, '[Leave / MC Request Form]',     'Digital submission and approval')
_scrshot(s, SCR_L2, SCR_T2, SCR_W, SCR_H, '[Attendance Export Report]',    'Excel and print-ready cycle summary')

# ============================================================
# 11.  EXPECTED BENEFITS AND IMPACT
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Expected Benefits and Impact'
tf = s.placeholders[12].text_frame; tf.word_wrap = True

_bullets(tf, [
    'Real-time presence visibility without being on site',
    'GPS verification prevents false reporting',
    'Full attendance history and trends for any cycle',
    'Early identification of absent-prone reservists',
    'Digital leave and MC handling end-to-end',
    'Permanent records replace memory and chat logs',
    'Reduced admin workload across all exercise cycles',
], sz=17, sp0=14, sp=16)

# ============================================================
# 12.  COST AND VALUE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Cost and Value'

# 4 sections filling to y=6.85 (no note), each row h=RH=2.77
_section_col(s, BX,          BY,     CW2, RH, 'OpsTracker Cost', GRN_PIL, [
    'Development: SGD 0 (built in-house)',
    'Hosting: SGD 0 (Netlify or GitHub Pages)',
    'Database: SGD 0 (Supabase free tier)',
    'Custom domain: ~SGD 15/year (optional)',
], sz=14)
_section_col(s, BX,          BOT_Y,  CW2, RH, 'Scale for This Unit', BLUE, [
    'At 5-7 reservists per week, under 1,000 records generated per year',
    'Free tier (50,000 rows) supports decades of operation at this scale',
    'Paid upgrade (SGD 25/month) unlikely to ever be needed',
    'Zero cost escalation risk for the foreseeable future',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,     CW2, RH, 'Commercial Equivalent', NAVY, [
    'Off-the-shelf HR SaaS: SGD 5-15 per user per month',
    '100 users: up to SGD 18,000 per year in licensing',
    'No commercial tool is built for NS exercise workflows',
    'OpsTracker: purpose-built at zero licensing cost',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y,  CW2, RH, 'Additional Value', AMBER, [
    'Source code fully owned by the unit',
    'Purpose-built for SPF NS shift and leave workflows',
    'No developer knowledge needed to operate',
    'Fully transferable to successor with zero cost',
], sz=14)

# ============================================================
# 13.  CHALLENGES
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Challenges'

CH = BH - 0.10  # full-height sections
_section_col(s, BX,              BY, CW2, CH, 'Technical', NAVY, [
    'GPS accuracy varies in covered or indoor environments',
    'In-app browsers restrict some PWA features',
    'Service worker versioning needed for reliable updates',
    'Realtime reconnection handling on unstable connections',
], sz=15)
_section_col(s, BX + CW2 + 0.20, BY, CW2, CH, 'Operational', BLUE, [
    'Shifting habits from verbal to app-based check-in',
    'Admin training required for full dashboard use',
    'Offline sync validation during field exercises',
    'Cross-device QA across iOS and Android',
], sz=15)

# ============================================================
# 14.  DEPLOYMENT TIMELINE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Deployment Timeline'

phases = [
    (NAVY,    'Approval',      'Aug 2026',           'Committee approval received; system configuration begins'),
    (BLUE,    'Admin Setup',   'Sep 2026',           'GPS boundaries set; first cycle created; admins trained'),
    (GRN_PIL, 'Soft Launch',   'Oct 2026',           'Deploy for the next Ops Security reservist cycle; validate in real conditions'),
    (GRN_PIL, 'Full Rollout',  'Nov 2026',           'All Ops Security reservists onboarded; live for all subsequent cycles'),
    (BLUE,    'Stabilise',     'Dec 2026 - Feb 2027','Full operations; feedback collected; refinements applied'),
    (AMBER,   'Handover Prep', 'Mar - Apr 2027',     'Docs finalised; successor designated and trained'),
    (NAVY,    'ORD',           'May 2027',           'Ownership transferred; NSPI Sun Jin ORDs'),
]

n_rows = len(phases) + 1
tbl = s.shapes.add_table(n_rows, 3, Inches(BX), Inches(BY), Inches(BW), Inches(BH)).table
for ci, w in enumerate([2.0, 2.2, 6.67]):
    tbl.columns[ci].width = Inches(w)

_cell(tbl.cell(0, 0), 'Phase',     bg=NAVY, fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 1), 'Timing',    bg=NAVY, fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 2), 'Milestone', bg=NAVY, fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)

for ri, (bg, phase, timing, milestone) in enumerate(phases):
    row_bg = LGRAY if ri % 2 == 0 else WHITE
    _cell(tbl.cell(ri+1, 0), phase,     bg=bg,     fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
    _cell(tbl.cell(ri+1, 1), timing,    bg=row_bg, fg=NAVY,  sz=12)
    _cell(tbl.cell(ri+1, 2), milestone, bg=row_bg, fg=NAVY,  sz=12)

# ============================================================
# 15.  HANDOVER PLAN
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Handover Plan'

# 4 sections filling to y=6.35 (note slide), each row h=RH_N=2.52
_section_col(s, BX,          BY,       CW2, RH_N, 'Role-Based Continuity', NAVY, [
    'No designated successor required in advance',
    'Whoever takes on the NS coordination role inherits the system',
    'Super Admin access transferred to incoming NSman or Regular Officer',
    'Onboarding time: approximately half a day with documentation',
], sz=14)
_section_col(s, BX,          BOT_Y_N,  CW2, RH_N, 'What Will Be Provided', BLUE, [
    'Admin user manual for all day-to-day operations',
    'SOP: cycle creation, onboarding, and report exports',
    'Troubleshooting guide for common issues',
    'Sun Jin contactable post-ORD for critical issues (best effort)',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,       CW2, RH_N, 'What Gets Transferred', NAVY, [
    'Supabase project ownership',
    'Hosting account (Netlify or GitHub Pages)',
    'Domain registrar credentials if applicable',
    'Source code repository access',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y_N,  CW2, RH_N, 'Sustainability by Design', GRN_PIL, [
    'All admin functions in-app; no command line needed',
    'Role-based access limits scope and risk',
    'System updates require no technical knowledge',
    'Free tier: no billing or renewal risk',
], sz=14)

_note(s, 'Designed so any incoming NSman stepping into the role can be fully operational within half a day.')

# ============================================================
# 16.  FUTURE PLANS AND RECOMMENDATION
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Future Plans and Recommendation'

COL_H = 2.50
_section_col(s, BX + 0*(CW3 + 0.20), BY, CW3, COL_H, 'Near-Term', NAVY, [
    'Push notifications for check-in windows',
    'Automated WhatsApp attendance reminders',
    'Multi-company support within one deployment',
], sz=13)
_section_col(s, BX + 1*(CW3 + 0.20), BY, CW3, COL_H, 'Mid-Term', BLUE, [
    'Attendance trend analytics and absence flags',
    'Native mobile app for stronger offline support',
    'Auto cycle reports submitted on close',
], sz=13)
_section_col(s, BX + 2*(CW3 + 0.20), BY, CW3, COL_H, 'Long-Term', GRN_PIL, [
    'SPF HR or NS portal integration',
    'Automated report submission to HQ',
    'Single sign-on with NS credentials',
], sz=13)

_divider(s, BY + COL_H + 0.18)

tf_r = _tf(s, BX, BY + COL_H + 0.30, BW, 3.10)
_p(tf_r,
   'OpsTracker is fully built, tested, and ready for deployment. We seek your approval to proceed before the next exercise cycle.',
   sz=16, bold=False, color=NAVY, sp=2, first=True)

for item in [
    'Deploy OpsTracker at the start of the next Ops Security reservist cycle',
    'Onboard all Ops Security reservists into the system',
    'Initiate IT security review for official and broader deployment',
]:
    p = tf_r.add_paragraph(); p.space_before = Pt(8)
    run = p.add_run(); run.text = '•  ' + item
    run.font.size = Pt(14); run.font.color.rgb = NAVY; run.font.name = 'Arial'

_p(tf_r,
   'Running cost: zero to minimal. With your approval, onboarding can begin immediately.',
   sz=15, bold=True, color=NAVY, sp=12)

# ============================================================
# SAVE
# ============================================================
prs.save('OpsTracker.pptx')
print(f'Done - {len(prs.slides)} slides saved to OpsTracker.pptx')
