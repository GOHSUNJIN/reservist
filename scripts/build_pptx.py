"""
OpsTracker - SPF NS Reservist Attendance Management
Approval deck. No em dashes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ---- PALETTE ----
NAVY    = RGBColor(0x00, 0x20, 0x5B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x4F, 0x81, 0xBD)
LGRAY   = RGBColor(0xD4, 0xE1, 0xF2)
CARD    = RGBColor(0xF2, 0xF6, 0xFB)   # very light card background
GREEN   = RGBColor(0x9B, 0xBB, 0x59)
RED_T   = RGBColor(0xFC, 0xE4, 0xD6)
GRN_T   = RGBColor(0xE2, 0xEF, 0xDA)
DK_RED  = RGBColor(0x7F, 0x1F, 0x1F)
DK_GRN  = RGBColor(0x1F, 0x4E, 0x1F)
GRN_PIL = RGBColor(0x26, 0x6B, 0x4A)
AMBER   = RGBColor(0xE0, 0x8A, 0x00)

AW = 0.07   # card left accent strip width (inches)

BX, BY, BW, BH = 1.78, 1.15, 10.87, 5.80
CW2 = (BW - 0.20) / 2
CW3 = (BW - 0.40) / 3
RH_N, RG_N = 2.52, 0.15
BOT_Y_N = BY + RH_N + RG_N
RH, RG = 2.77, 0.15
BOT_Y = BY + RH + RG

# ---- HELPERS ----

def remove_all_slides(prs):
    ids = prs.slides._sldIdLst
    for elem in list(ids):
        rId = elem.get(qn('r:id'))
        ids.remove(elem)
        prs.part.drop_rel(rId)

def _tf(slide, l, t, w, h):
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    sh.text_frame.word_wrap = True
    return sh.text_frame

def _p(tf, text, sz=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
       sp=6, first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(sp)
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = 'Arial'
    return p

def _pill(slide, l, t, w, h, label, bg, fg=WHITE, sz=13):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = bg; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(sz); run.font.bold = True
    run.font.color.rgb = fg; run.font.name = 'Arial'

def _bullets(tf, items, sz=14, sp0=8, sp=11, color=NAVY):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(sp0 if i == 0 else sp)
        run = p.add_run(); run.text = '•  ' + item
        run.font.size = Pt(sz); run.font.color.rgb = color; run.font.name = 'Arial'

def _section_col(slide, l, t, w, h, label, bg, items, sz=14):
    """Card with left accent strip, colored label, thin rule, then bullets."""
    card = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = CARD; card.line.fill.background()
    acc = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(AW), Inches(h))
    acc.fill.solid(); acc.fill.fore_color.rgb = bg; acc.line.fill.background()
    tf_lbl = _tf(slide, l + AW + 0.12, t + 0.11, w - AW - 0.18, 0.30)
    _p(tf_lbl, label, sz=11, bold=True, color=bg, first=True)
    rule = slide.shapes.add_shape(1, Inches(l + AW + 0.12), Inches(t + 0.43),
                                   Inches(w - AW - 0.20), Inches(0.016))
    rule.fill.solid(); rule.fill.fore_color.rgb = bg; rule.line.fill.background()
    tf = _tf(slide, l + AW + 0.12, t + 0.48, w - AW - 0.20, h - 0.56)
    _bullets(tf, items, sz=sz, sp0=4, sp=9, color=NAVY)

def _flow_step(slide, l, t, w, h, step_label, bg, action, details):
    """Flow card with left accent, step badge, action title, rule, details."""
    card = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = CARD; card.line.fill.background()
    acc = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(AW), Inches(h))
    acc.fill.solid(); acc.fill.fore_color.rgb = bg; acc.line.fill.background()
    _pill(slide, l + AW + 0.10, t + 0.12, w - AW - 0.20, 0.28, step_label, bg, sz=10)
    tf_act = _tf(slide, l + AW + 0.10, t + 0.50, w - AW - 0.20, 0.60)
    _p(tf_act, action, sz=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER, first=True)
    rule = slide.shapes.add_shape(1, Inches(l + AW + 0.12), Inches(t + 1.14),
                                   Inches(w - AW - 0.22), Inches(0.016))
    rule.fill.solid(); rule.fill.fore_color.rgb = bg; rule.line.fill.background()
    tf_det = _tf(slide, l + AW + 0.12, t + 1.22, w - AW - 0.22, h - 1.30)
    _bullets(tf_det, details, sz=13, sp0=6, sp=10, color=NAVY)

def _row_card(slide, l, t, w, h, label, bg, desc, sz_lbl=11, sz_desc=13):
    """Single-row card with left accent, label, rule, description. For lists."""
    card = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = CARD; card.line.fill.background()
    acc = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(AW), Inches(h))
    acc.fill.solid(); acc.fill.fore_color.rgb = bg; acc.line.fill.background()
    tf_lbl = _tf(slide, l + AW + 0.12, t + 0.09, w - AW - 0.18, 0.24)
    _p(tf_lbl, label, sz=sz_lbl, bold=True, color=bg, first=True)
    rule = slide.shapes.add_shape(1, Inches(l + AW + 0.12), Inches(t + 0.35),
                                   Inches(w - AW - 0.20), Inches(0.013))
    rule.fill.solid(); rule.fill.fore_color.rgb = bg; rule.line.fill.background()
    tf_desc = _tf(slide, l + AW + 0.12, t + 0.38, w - AW - 0.20, h - 0.44)
    _p(tf_desc, desc, sz=sz_desc, color=NAVY, first=True, sp=0)

def _cell(cell, text, bg=None, fg=NAVY, sz=12, bold=False, align=PP_ALIGN.LEFT):
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell._tc.get_or_add_tcPr().set('anchor', 'ctr')
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for r in list(p.runs): p._p.remove(r._r)
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.color.rgb = fg; run.font.name = 'Arial'

def _note(slide, text, y=6.45, sz=10):
    tf = _tf(slide, BX, y, BW, 0.42)
    _p(tf, text, sz=sz, color=BLUE, italic=True, first=True)

def _pill_row(slide, items, t, h=0.44):
    n = len(items)
    pw = (BW - 0.15 * (n - 1)) / n
    for i, (lbl, bg) in enumerate(items):
        _pill(slide, BX + i * (pw + 0.15), t, pw, h, lbl, bg)

def _section_slide(title):
    s = prs.slides.add_slide(L0)
    ph = s.placeholders[0]; ph.text = ''
    p = ph.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(36); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = 'Arial'
    for idx in [11, 12, 13]:
        try: s.placeholders[idx].text = ''
        except Exception: pass
    return s

# ---- LOAD ----
prs = Presentation('OpsTracker.pptx')
L0 = prs.slide_layouts[0]
L3 = prs.slide_layouts[3]
remove_all_slides(prs)

# ============================================================
# 1.  TITLE
# ============================================================
s = prs.slides.add_slide(L0)
ph = s.placeholders[0]; ph.text = ''
tf0 = ph.text_frame
pm = tf0.paragraphs[0]; pm.alignment = PP_ALIGN.CENTER
rm = pm.add_run(); rm.text = 'OpsTracker'
rm.font.size = Pt(44); rm.font.bold = True
rm.font.color.rgb = WHITE; rm.font.name = 'Arial'
ps = tf0.add_paragraph(); ps.alignment = PP_ALIGN.CENTER; ps.space_before = Pt(8)
rs = ps.add_run(); rs.text = 'Digital Attendance Management for SPF NS Reservists'
rs.font.size = Pt(20); rs.font.color.rgb = WHITE; rs.font.name = 'Arial'
s.placeholders[11].text = '17 August 2026'
s.placeholders[12].text = 'For Approval'
tf13 = s.placeholders[13].text_frame
p1 = tf13.paragraphs[0]; r1 = p1.add_run()
r1.text = 'Ops Security, Ops Branch'; r1.font.bold = True; r1.font.name = 'Arial'
p2 = tf13.add_paragraph(); r2 = p2.add_run()
r2.text = 'NSPI Sun Jin'; r2.font.name = 'Arial'

# ============================================================
# 2.  AGENDA
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Agenda'

agenda = [
    ('01', 'Executive Summary'),
    ('02', 'The Problem'),
    ('03', 'Before vs After'),
    ('04', 'Introducing OpsTracker'),
    ('05', 'Reservist Experience'),
    ('06', 'Admin Capabilities and Features'),
    ('07', 'Reservist App Screens'),
    ('08', 'Admin Dashboard Screens'),
    ('09', 'Expected Benefits and Impact'),
    ('10', 'Data Security and Privacy'),
    ('11', 'PDPA Compliance'),
    ('12', 'Cost and Value'),
    ('13', 'Why Now'),
    ('14', 'Challenges'),
    ('15', 'Deployment Timeline'),
    ('16', 'Handover Plan'),
    ('17', 'Future Plans and Recommendation'),
]

ITEM_H = 0.50
PILL_W = 0.34
PILL_H = 0.25

def _agenda_col(slide, l, rows, accent):
    for i, (num, label) in enumerate(rows):
        ty = BY + i * ITEM_H
        _pill(slide, l, ty + 0.02, PILL_W, PILL_H, num, accent, sz=9)
        div = slide.shapes.add_shape(1, Inches(l + PILL_W + 0.06), Inches(ty + 0.06),
                                      Inches(0.016), Inches(0.20))
        div.fill.solid(); div.fill.fore_color.rgb = LGRAY; div.line.fill.background()
        tf = _tf(slide, l + PILL_W + 0.12, ty + 0.01, 4.60, 0.34)
        _p(tf, label, sz=12, color=NAVY, first=True, sp=2)

_agenda_col(s, BX,        agenda[:9],  NAVY)
_agenda_col(s, BX + 5.55, agenda[9:],  BLUE)

# ============================================================
# 3.  EXECUTIVE SUMMARY
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Executive Summary'

_section_col(s, BX,               BY,      CW2, RH_N, 'The Challenge', NAVY, [
    'Attendance tracked verbally with no system verification',
    'No real-time visibility for the team',
    'No consolidated records or audit trail',
], sz=15)
_section_col(s, BX + CW2 + 0.20,  BY,      CW2, RH_N, 'The Solution', GRN_PIL, [
    'GPS-verified check-in via mobile browser',
    'Live dashboard viewable by the team',
    'Digital leave and MC; full history; one-click export',
], sz=15)
_section_col(s, BX,               BOT_Y_N, CW2, RH_N, 'Cost to the Unit', BLUE, [
    'Built in-house at SGD 0 development cost',
    'Running cost: SGD 0 to SGD 15 per year',
    'Free tier supports years of use at this unit scale',
], sz=15)
_section_col(s, BX + CW2 + 0.20,  BOT_Y_N, CW2, RH_N, 'What We Are Asking For', AMBER, [
    'Approval to deploy at the next Ops Security cycle',
    'System is fully built, tested, and ready to go',
    'Onboarding can begin within days of approval',
], sz=15)

# ============================================================
# 4.  DIVIDER: THE PROBLEM
# ============================================================
_section_slide('The Problem')

# ============================================================
# 5.  THE PROBLEM
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'The Problem'

ch = (BH - 0.20) / 3
cg = 0.10

_section_col(s, BX, BY,             BW, ch, '01  No Verification', NAVY, [
    'Reservists report verbally with no system check or confirmation',
    'No way to confirm physical presence at the reporting location',
], sz=15)
_section_col(s, BX, BY + ch + cg,   BW, ch, '02  No Visibility',   BLUE, [
    'The team has no real-time view without being physically present',
    'Leave and MC handled informally via chat with no formal record',
], sz=15)
_section_col(s, BX, BY + 2*(ch+cg), BW, ch, '03  No Records',      NAVY, [
    'No consolidated attendance history exists for the unit',
    'Chronic absences go undetected until it is too late to act',
], sz=15)

# ============================================================
# 6.  BEFORE vs AFTER
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Before vs After'

bva = [
    ('Attendance Check', 'Verbal, trust-based; no verification',
     'GPS-verified check-in via mobile browser'),
    ('Supervisor View',  'Must be on site to know who showed up',
     'Live dashboard updated in real time'),
    ('Leave and MC',     'WhatsApp or verbal; no formal record',
     'Digital submission with full approval history'),
    ('Record Keeping',   'Relies on memory and scattered chat logs',
     'Permanent digital log with one-click export'),
    ('Accountability',   'No mechanism to flag repeated no-shows',
     'Auto-absent marking and trend tracking'),
]

tbl = s.shapes.add_table(6, 3, Inches(BX - 0.18), Inches(BY), Inches(BW), Inches(BH)).table
for idx, w in enumerate([2.2, 4.15, 4.52]):
    tbl.columns[idx].width = Inches(w)
_cell(tbl.cell(0, 0), 'Area',       bg=NAVY, fg=WHITE, sz=13, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 1), 'Before',     bg=NAVY, fg=WHITE, sz=13, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 2), 'OpsTracker', bg=BLUE, fg=WHITE, sz=13, bold=True, align=PP_ALIGN.CENTER)
for ri, (area, before, after) in enumerate(bva):
    alt = ri % 2 == 0
    _cell(tbl.cell(ri+1, 0), area,   bg=LGRAY if alt else WHITE, fg=NAVY, sz=13, bold=True)
    _cell(tbl.cell(ri+1, 1), before, bg=RED_T, fg=DK_RED, sz=13)
    _cell(tbl.cell(ri+1, 2), after,  bg=GRN_T, fg=DK_GRN, sz=13)

# ============================================================
# 7.  DIVIDER: THE SOLUTION
# ============================================================
_section_slide('The Solution')

# ============================================================
# 8.  INTRODUCING OPSTRACKER
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Introducing OpsTracker'

tf_i = _tf(s, BX, BY, BW, 0.42)
_p(tf_i,
   'Mobile-first Progressive Web App replacing verbal check-ins with GPS-verified digital attendance records.',
   sz=15, color=NAVY, first=True)

_pill_row(s, [
    ('Progressive Web App',    NAVY),
    ('Supabase + Realtime',    BLUE),
    ('GPS Shift Verification', GRN_PIL),
], t=1.68)

features = [
    ('Mobile-First',        'Works on iOS, Android, or desktop. No app store download required.'),
    ('GPS Check-In',        'Confirms device location before recording. Supports AM, PM, and Full Day shifts.'),
    ('Real-Time Dashboard', 'Live view of present, absent, and MC status as check-ins happen.'),
    ('Attendance History',  'Full records and trends for any reservist or cycle, searchable anytime.'),
    ('Automated Reporting', 'One-click export to Excel or a print-ready cycle report.'),
]

FEAT_H = 0.76
FEAT_G = 0.09
feat_colors = [NAVY, BLUE, GRN_PIL, BLUE, NAVY]
for i, ((lbl, det), col) in enumerate(zip(features, feat_colors)):
    fy = 2.24 + i * (FEAT_H + FEAT_G)
    _row_card(s, BX, fy, BW, FEAT_H, lbl, col, det, sz_lbl=13, sz_desc=13)

# ============================================================
# 9.  RESERVIST EXPERIENCE  (horizontal flow)
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Reservist Experience'

FA = 0.38
FW = (BW - 3 * FA) / 4
FH = 4.80

flow = [
    ('Step 1', NAVY,    'Register',
     ['Submit form online', 'Admin activates account', 'Receive login credentials']),
    ('Step 2', BLUE,    'Check In',
     ['Open app in reporting window', 'GPS confirms location', 'One tap records attendance']),
    ('Step 3', GRN_PIL, 'Leave / MC',
     ['Submit request digitally', 'Receive approval in-app', 'History tracked automatically']),
    ('Step 4', GREEN,   'View History',
     ['Access records anytime', 'Check attendance trends', 'See leave outcomes']),
]

for i, (step, bg, action, details) in enumerate(flow):
    lx = BX + i * (FW + FA)
    _flow_step(s, lx, BY, FW, FH, step, bg, action, details)
    if i < 3:
        tf_a = _tf(s, lx + FW + 0.04, BY + FH / 2 - 0.22, FA - 0.08, 0.44)
        _p(tf_a, '▶', sz=18, bold=False, color=BLUE, align=PP_ALIGN.CENTER, first=True)

_note(s, 'Reservists can install OpsTracker to their home screen as a PWA. No app store required.')

# ============================================================
# 10.  ADMIN AND SUPERVISOR CAPABILITIES
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Admin and Supervisor Capabilities'

_section_col(s, BX,               BY,      CW2, RH_N, 'Attendance', NAVY, [
    'Live dashboard: present, absent, and MC',
    'Manual attendance override with full audit log',
    'Mark all present or configure no-report days',
    'Auto-absent marking after check-in window closes',
], sz=14)
_section_col(s, BX,               BOT_Y_N, CW2, RH_N, 'Leave and MC', BLUE, [
    'Approve or reject requests individually or in bulk',
    'Cancel approved leave if plans change',
    'Full leave history per reservist',
    'MC tracked separately from standard leave',
], sz=14)
_section_col(s, BX + CW2 + 0.20,  BY,      CW2, RH_N, 'Roster', NAVY, [
    'Add, deactivate, or reassign reservists',
    'CSV bulk upload and password resets',
    'Welfare notes and member activity log',
    'Assign reservists to batches and companies',
], sz=14)
_section_col(s, BX + CW2 + 0.20,  BOT_Y_N, CW2, RH_N, 'Reporting', BLUE, [
    'Export cycle report to Excel or print',
    'WhatsApp broadcast to the full unit',
    'Cycle picker and multi-day date range view',
    'Attendance trend snapshot per exercise cycle',
], sz=14)
_note(s, 'Super Admins can also manage admin accounts, create cycles, and define GPS reporting boundaries.')

# ============================================================
# 11.  RESERVIST APP SCREENS
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Reservist App Screens'

SCR_W     = (BW - 3 * 0.15) / 4   # 2.605" per screen
SCR_H     = SCR_W * 16 / 9         # 4.631" portrait
SCR_LBL_Y = BY + SCR_H + 0.14      # label row y

def _scrshot_p(slide, i, placeholder, feature, desc):
    l = BX + i * (SCR_W + 0.15)
    # Phone frame (NAVY outer)
    frame = slide.shapes.add_shape(1, Inches(l), Inches(BY - 0.04),
                                    Inches(SCR_W), Inches(SCR_H + 0.04))
    frame.fill.solid(); frame.fill.fore_color.rgb = NAVY; frame.line.fill.background()
    # White screen area
    screen = slide.shapes.add_shape(1, Inches(l), Inches(BY), Inches(SCR_W), Inches(SCR_H))
    screen.fill.solid(); screen.fill.fore_color.rgb = WHITE; screen.line.fill.background()
    # Status bar
    sbar = slide.shapes.add_shape(1, Inches(l), Inches(BY), Inches(SCR_W), Inches(0.22))
    sbar.fill.solid(); sbar.fill.fore_color.rgb = NAVY; sbar.line.fill.background()
    # Placeholder label inside screen
    tf = _tf(slide, l + 0.10, BY + SCR_H * 0.38, SCR_W - 0.20, SCR_H * 0.28)
    _p(tf, placeholder, sz=10, color=LGRAY, align=PP_ALIGN.CENTER, first=True)
    # Feature label below
    tf2 = _tf(slide, l, SCR_LBL_Y, SCR_W, 0.54)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = feature
    r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = NAVY; r2.font.name = 'Arial'
    p3 = tf2.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(3)
    r3 = p3.add_run(); r3.text = desc
    r3.font.size = Pt(9); r3.font.italic = True; r3.font.color.rgb = BLUE; r3.font.name = 'Arial'

_scrshot_p(s, 0, '[Login Screen]',
           'Login', 'Email and password sign-in; session persists on device')
_scrshot_p(s, 1, '[Check-In Screen]',
           'GPS Check-In', 'Shift selector, GPS confirmation, one-tap submission')
_scrshot_p(s, 2, '[Leave / MC Form]',
           'Leave and MC', 'Date range, reason, and type; submitted digitally')
_scrshot_p(s, 3, '[My Attendance]',
           'Attendance History', 'Personal record of all shifts, leaves, and MC')

# ============================================================
# 12.  ADMIN DASHBOARD SCREENS
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Admin Dashboard Screens'

_scrshot_p(s, 0, '[Attendance Dashboard]',
           'Live Dashboard', 'Real-time present, absent, and MC status per shift')
_scrshot_p(s, 1, '[Leave Approvals Queue]',
           'Leave Approvals', 'Pending requests with one-tap approve or reject')
_scrshot_p(s, 2, '[Roster Management]',
           'Roster', 'Reservist list, status, welfare notes, and batch assignment')
_scrshot_p(s, 3, '[Cycle Report Export]',
           'Reports and Export', 'One-click Excel export or print-ready cycle report')

# ============================================================
# 13.  EXPECTED BENEFITS AND IMPACT
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Expected Benefits and Impact'

BRH = 0.74   # benefit row card height
BRG = 0.07   # gap between rows

bene_rows = [
    ('Real-Time Visibility',    BLUE,
     'See who is present, absent, or on MC without being physically on site'),
    ('GPS Verification',        GRN_PIL,
     'Device location confirmed before each check-in is recorded'),
    ('Full Attendance History',  NAVY,
     'Complete attendance records and trends for every reservist and cycle'),
    ('Early Accountability',    AMBER,
     'Absent patterns identified automatically before they become critical'),
    ('Digital Leave and MC',    BLUE,
     'End-to-end digital handling with a full approval and rejection history'),
    ('Permanent Records',       NAVY,
     'Replaces memory and scattered chat logs with a verified digital trail'),
    ('Reduced Admin Work',      GRN_PIL,
     'One-click exports and live dashboards cut manual coordination effort'),
]

for i, (lbl, col, desc) in enumerate(bene_rows):
    _row_card(s, BX, BY + i * (BRH + BRG), BW, BRH, lbl, col, desc)

# ============================================================
# 14.  DIVIDER: SECURITY AND COMPLIANCE
# ============================================================
_section_slide('Security and Compliance')

# ============================================================
# 15.  DATA SECURITY AND PRIVACY
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Data Security and Privacy'

_section_col(s, BX,              BY,     CW2, RH, 'Access Control', NAVY, [
    'Three-tier roles: Reservist, Admin, Super Admin',
    'Each role restricted to authorised data only',
    'Super Admin approval required for account changes',
    'All actions tied to authenticated identity',
], sz=14)
_section_col(s, BX,              BOT_Y,  CW2, RH, 'Session Security', BLUE, [
    'Automatic logout after a period of inactivity',
    'Secure auth tokens managed by Supabase Auth',
    'No personal data stored on the device',
    'Login required on every new session',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,     CW2, RH, 'Data Protection', NAVY, [
    'PostgreSQL: encrypted at rest and in transit (TLS)',
    'Row-level security enforced at database layer',
    'No third-party data sharing',
    'ISO 27001 certified; SOC 2 Type II compliant',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y,  CW2, RH, 'Audit Trail', GRN_PIL, [
    'Every edit logged with timestamp and user',
    'Leave decisions and admin actions recorded',
    'Full history available for compliance review',
    'Immutable log; edits cannot be silently deleted',
], sz=14)
_note(s, 'Hosted externally on Supabase - comparable in approach to Google Forms, with stronger security controls. No SPF infrastructure required.')

# ============================================================
# 16.  PDPA COMPLIANCE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'PDPA Compliance'

tf_i = _tf(s, BX, BY, BW, 0.44)
_p(tf_i,
   'Designed to align with the Personal Data Protection Act 2012 (PDPA) and government data governance principles.',
   sz=14, color=NAVY, first=True)

P_TOP = BY + 0.60
P_RH  = 2.43
P_BOT = P_TOP + P_RH + 0.17

_section_col(s, BX,              P_TOP, CW2, P_RH, 'Data Collected', NAVY, [
    'Name and contact number (at registration)',
    'GPS location (at check-in only)',
    'Attendance and leave records',
    'No NRIC or financial data collected',
], sz=14)
_section_col(s, BX,              P_BOT, CW2, P_RH, 'Individual Rights', BLUE, [
    'Reservists may request their own records',
    'Corrections made by authorised admin',
    'Data deleted after NS liability ends',
    'Requests handled within a reasonable time',
], sz=14)
_section_col(s, BX + CW2 + 0.20, P_TOP, CW2, P_RH, 'Key Obligations Met', NAVY, [
    'Consent obtained at sign-up before data collected',
    'Purpose limited to attendance management only',
    'Only necessary data collected (data minimisation)',
    'Encryption, access control, row-level security',
], sz=14)
_section_col(s, BX + CW2 + 0.20, P_BOT, CW2, P_RH, 'Recommended Actions', AMBER, [
    'Seek Data Protection Officer guidance before deployment',
    'Define data retention period per SPF policy',
    'Include a privacy notice at reservist sign-up',
    'Document the lawful basis for data collection',
], sz=14)

# ============================================================
# 17.  DIVIDER: THE BUSINESS CASE
# ============================================================
_section_slide('The Business Case')

# ============================================================
# 18.  COST AND VALUE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Cost and Value'

_section_col(s, BX,              BY,     CW2, RH, 'OpsTracker Cost', GRN_PIL, [
    'Development: SGD 0 (built in-house)',
    'Hosting: SGD 0 (Netlify or GitHub Pages)',
    'Database: SGD 0 (Supabase free tier)',
    'Custom domain: ~SGD 15/year (optional)',
], sz=14)
_section_col(s, BX,              BOT_Y,  CW2, RH, 'Scale for This Unit', BLUE, [
    'At 5 to 7 reservists per week, under 1,000 records per year',
    'Free tier (50,000 rows) supports decades at this scale',
    'Paid upgrade (SGD 25/month) unlikely to ever be needed',
    'Zero cost escalation risk for the foreseeable future',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,     CW2, RH, 'Commercial Equivalent', NAVY, [
    'Off-the-shelf HR SaaS: SGD 5 to 15 per user per month',
    '100 users: up to SGD 18,000 per year in licensing',
    'No commercial tool is built for NS exercise workflows',
    'OpsTracker: purpose-built at zero licensing cost',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y,  CW2, RH, 'Additional Value', AMBER, [
    'Source code fully owned by the unit',
    'Purpose-built for SPF NS shift and leave workflows',
    'No developer knowledge needed to operate',
    'Fully transferable to successor at zero cost',
], sz=14)

# ============================================================
# 19.  WHY NOW
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Why Now'

cww = (BW - 0.30) / 3

_section_col(s, BX + 0 * (cww + 0.15), BY, cww, BH - 0.10, 'Next Cycle Is the Window', BLUE, [
    'The next Ops Security cycle is the optimal deployment window',
    'Delaying means another cycle without records or accountability',
    'Admin setup and onboarding require 4 to 6 weeks lead time',
    'Starting now ensures the system is ready when reservists report',
], sz=14)
_section_col(s, BX + 1 * (cww + 0.15), BY, cww, BH - 0.10, 'ORD Deadline: May 2027', AMBER, [
    'NSPI Sun Jin ORDs in May 2027',
    'System must be stable and handed over by March 2027',
    'Deployment now gives 6 months of stabilisation time',
    'A delayed start risks a rushed or incomplete handover',
], sz=14)
_section_col(s, BX + 2 * (cww + 0.15), BY, cww, BH - 0.10, 'Zero Cost to Approve', GRN_PIL, [
    'OpsTracker is already built and tested; no further dev cost',
    'Running cost: SGD 0 to SGD 15 per year at this scale',
    'Every cycle without it is a cycle without accountability',
    'System is ready the moment approval is given',
], sz=14)

# ============================================================
# 20.  CHALLENGES
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Challenges'

CH = BH - 0.10
_section_col(s, BX,              BY, CW2, CH, 'Technical', NAVY, [
    'GPS accuracy can vary indoors or in covered spaces',
    'Some in-app browsers limit PWA features',
    'Service worker versioning needed for reliable updates',
    'Network reconnection handling in poor signal areas',
], sz=15)
_section_col(s, BX + CW2 + 0.20, BY, CW2, CH, 'Operational', BLUE, [
    'Encouraging reservists to check in via the app',
    'Admin training required for full dashboard use',
    'Sync reliability to be validated before field use',
    'Cross-device QA across iOS and Android',
], sz=15)

# ============================================================
# 21.  DEPLOYMENT TIMELINE
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Deployment Timeline'

phases = [
    (NAVY,    'Approval',      'Aug 2026',            'Committee approval received; system configuration begins'),
    (BLUE,    'Admin Setup',   'Sep 2026',            'GPS boundaries set; first cycle created; admins trained'),
    (GRN_PIL, 'Soft Launch',   'Oct 2026',            'Deploy for the next Ops Security reservist cycle; validate in real conditions'),
    (GRN_PIL, 'Full Rollout',  'Nov 2026',            'All Ops Security reservists onboarded; live for all subsequent cycles'),
    (BLUE,    'Stabilise',     'Dec 2026 - Feb 2027', 'Full operations; feedback collected; refinements applied'),
    (AMBER,   'Handover Prep', 'Mar - Apr 2027',      'Docs finalised; successor designated and trained'),
    (NAVY,    'ORD',           'May 2027',            'Ownership transferred; NSPI Sun Jin ORDs'),
]

tbl = s.shapes.add_table(8, 3, Inches(BX), Inches(BY), Inches(BW), Inches(BH)).table
for ci, w in enumerate([2.0, 2.2, 6.67]):
    tbl.columns[ci].width = Inches(w)
_cell(tbl.cell(0, 0), 'Phase',     bg=NAVY, fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 1), 'Timing',    bg=NAVY, fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
_cell(tbl.cell(0, 2), 'Milestone', bg=NAVY, fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
for ri, (bg, phase, timing, milestone) in enumerate(phases):
    row_bg = CARD if ri % 2 == 0 else WHITE
    _cell(tbl.cell(ri+1, 0), phase,     bg=bg,     fg=WHITE, sz=12, bold=True, align=PP_ALIGN.CENTER)
    _cell(tbl.cell(ri+1, 1), timing,    bg=row_bg, fg=NAVY,  sz=12)
    _cell(tbl.cell(ri+1, 2), milestone, bg=row_bg, fg=NAVY,  sz=12)

# ============================================================
# 22.  HANDOVER PLAN
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Handover Plan'

_section_col(s, BX,              BY,      CW2, RH_N, 'Role-Based Continuity', NAVY, [
    'No designated successor required in advance',
    'Whoever takes the NS coordination role inherits the system',
    'Super Admin access transferred to incoming NSman or Regular',
    'Onboarding time: approximately half a day with documentation',
], sz=14)
_section_col(s, BX,              BOT_Y_N, CW2, RH_N, 'What Will Be Provided', BLUE, [
    'Admin user manual for all day-to-day operations',
    'SOP: cycle creation, onboarding, and report exports',
    'Troubleshooting guide for common issues',
    'Sun Jin contactable post-ORD for critical issues (best effort)',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BY,      CW2, RH_N, 'What Gets Transferred', NAVY, [
    'Supabase project ownership',
    'Hosting account (Netlify or GitHub Pages)',
    'Domain registrar credentials if applicable',
    'Source code repository access',
], sz=14)
_section_col(s, BX + CW2 + 0.20, BOT_Y_N, CW2, RH_N, 'Sustainability by Design', GRN_PIL, [
    'All admin functions in-app; no command line needed',
    'Role-based access limits scope and risk',
    'System updates require no technical knowledge',
    'Free tier: no billing or renewal risk',
], sz=14)
_note(s, 'Designed so any incoming NSman stepping into the role can be fully operational within half a day.')

# ============================================================
# 23.  FUTURE PLANS AND RECOMMENDATION
# ============================================================
s = prs.slides.add_slide(L3)
s.placeholders[0].text = 'Future Plans and Recommendation'

COL_H = 2.40
_section_col(s, BX + 0*(CW3 + 0.20), BY, CW3, COL_H, 'Near-Term', NAVY, [
    'Push notifications for check-in windows',
    'Automated WhatsApp attendance reminders',
    'Multi-company support within a single cycle',
], sz=13)
_section_col(s, BX + 1*(CW3 + 0.20), BY, CW3, COL_H, 'Mid-Term', BLUE, [
    'Attendance trend analytics and absence alerts',
    'Native mobile app with stronger offline support',
    'Auto-generate and send cycle reports on close',
], sz=13)
_section_col(s, BX + 2*(CW3 + 0.20), BY, CW3, COL_H, 'Long-Term', GRN_PIL, [
    'SPF HR or NS portal integration',
    'Automated report submission to HQ',
    'Single sign-on with NS credentials',
], sz=13)

# Recommendation box
rec_t = BY + COL_H + 0.22
rec_h = BH - COL_H - 0.28
rec = s.shapes.add_shape(1, Inches(BX), Inches(rec_t), Inches(BW), Inches(rec_h))
rec.fill.solid(); rec.fill.fore_color.rgb = NAVY; rec.line.fill.background()
# Top accent stripe on rec box
rec_stripe = s.shapes.add_shape(1, Inches(BX), Inches(rec_t), Inches(BW), Inches(0.055))
rec_stripe.fill.solid(); rec_stripe.fill.fore_color.rgb = BLUE; rec_stripe.line.fill.background()

tf_r = rec.text_frame; tf_r.word_wrap = True

p = tf_r.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(16)
r = p.add_run(); r.text = 'We Are Asking For Your Approval'
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = 'Arial'

for item in [
    'Deploy OpsTracker at the start of the next Ops Security reservist cycle',
    'Onboard all Ops Security reservists into the system',
    'Initiate IT security review for official and broader deployment',
]:
    p = tf_r.add_paragraph(); p.space_before = Pt(8)
    r = p.add_run(); r.text = '    •  ' + item
    r.font.size = Pt(14); r.font.color.rgb = WHITE; r.font.name = 'Arial'

p = tf_r.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(14)
r = p.add_run()
r.text = 'System is ready now.  Running cost: SGD 0 to SGD 15 per year.  Onboarding can begin immediately.'
r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = 'Arial'

# ============================================================
# SAVE
# ============================================================
prs.save('OpsTracker.pptx')
print(f'Done - {len(prs.slides)} slides saved to OpsTracker.pptx')
